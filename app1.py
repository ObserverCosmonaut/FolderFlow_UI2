import streamlit as st
import tempfile
import os
from typing import List, Dict
from pathlib import Path
import base64
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.schema import Document
import google.generativeai as genai
from pptx import Presentation

class CustomPPTLoader:
    def __init__(self, file_path: str, original_filename: str):
        self.file_path = file_path
        self.original_filename = original_filename

    def load(self) -> List[Document]:
        try:
            prs = Presentation(self.file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                if text_content:
                    text = "\n\n".join(text_content)
                    metadata = {"source": self.original_filename, "slide_number": slide_num}
                    documents.append(Document(page_content=text, metadata=metadata))
            
            return documents
        except Exception as e:
            st.error(f"Error processing PowerPoint file: {str(e)}")
            return []

class DocumentManager:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self.processed_files = []
        self.qa_chain = None
        self.chat_history = []
        self.total_files = 0
        self.processed_count = 0

    def process_file(self, uploaded_file) -> List[Document]:
        try:
            original_filename = uploaded_file.name
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(original_filename)[1]) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                file_path = tmp_file.name

            if original_filename.endswith('.pdf'):
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                for doc in documents:
                    doc.metadata['source'] = original_filename
            elif original_filename.endswith('.docx'):
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                for doc in documents:
                    doc.metadata['source'] = original_filename
            elif original_filename.endswith(('.pptx', '.ppt')):
                loader = CustomPPTLoader(file_path, original_filename)
                documents = loader.load()
            else:
                return []
            
            if documents:
                chunks = self.text_splitter.split_documents(documents)
                self.processed_files.append(original_filename)
                self.processed_count += 1
                return chunks
            return []

        except Exception as e:
            st.error(f"Error processing {original_filename}: {str(e)}")
            return []
        finally:
            if 'file_path' in locals():
                try:
                    os.unlink(file_path)
                except:
                    pass

    def process_folder(self, folder_files: List) -> List[Document]:
        all_chunks = []
        for file in folder_files:
            if file.name.lower().endswith(('.pdf', '.docx', '.pptx', '.ppt')):
                chunks = self.process_file(file)
                all_chunks.extend(chunks)
        return all_chunks

    def setup_qa_system(self, uploaded_files, is_folder=False):
        try:
            all_chunks = []
            self.total_files = len(uploaded_files)
            self.processed_count = 0
            
            progress_bar = st.progress(0)
            status_text = st.empty()

            if is_folder:
                status_text.text("Processing folder contents...")
                all_chunks = self.process_folder(uploaded_files)
            else:
                for uploaded_file in uploaded_files:
                    status_text.text(f"Processing {uploaded_file.name}...")
                    chunks = self.process_file(uploaded_file)
                    all_chunks.extend(chunks)
                    progress = int((self.processed_count / self.total_files) * 100)
                    progress_bar.progress(progress)
                    status_text.text(f"Processed {self.processed_count} of {self.total_files} files ({progress}%)")

            if not all_chunks:
                st.error("No documents were successfully processed!")
                return False

            status_text.text("Setting up QA system...")
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            vectorstore = FAISS.from_documents(all_chunks, embeddings)
            
            self.qa_chain = ConversationalRetrievalChain.from_llm(
                ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.7),
                vectorstore.as_retriever(),
                return_source_documents=True
            )
            
            status_text.text("System ready!")
            progress_bar.progress(100)
            return True

        except Exception as e:
            st.error(f"Error setting up QA system: {str(e)}")
            return False

    def ask_question(self, question: str) -> Dict:
        if not self.qa_chain:
            return {"error": "QA system not initialized. Please upload documents first."}
        try:
            result = self.qa_chain.invoke({
                "question": question,
                "chat_history": self.chat_history
            })
            
            sources = list({
                doc.metadata.get('source', 'Unknown source')
                for doc in result["source_documents"]
            })
            
            self.chat_history.append((question, result["answer"]))
            return {
                "answer": result["answer"],
                "sources": sources
            }
        except Exception as e:
            return {"error": f"Error processing question: {str(e)}"}

def add_folder_upload_button():
    folder_upload_html = """
        <input type="file" id="folder-upload" webkitdirectory directory multiple style="display: none;">
        <label for="folder-upload" class="folder-upload-label">
            Upload Folder
        </label>
        
        <style>
            .folder-upload-label {
                display: inline-block;
                padding: 0.5rem 1rem;
                background-color: #D00F22;
                color: white;
                border-radius: 0.3rem;
                cursor: pointer;
                margin: 0.5rem 0;
                text-align: center;
                width: 100%;
                box-sizing: border-box;
            }
            .folder-upload-label:hover {
                background-color: #B00D1E;
            }
        </style>
        
        <script>
            const folderInput = document.getElementById('folder-upload');
            folderInput.addEventListener('change', function(e) {
                const files = Array.from(e.target.files);
                const fileData = files.map(file => ({
                    name: file.name,
                    path: file.webkitRelativePath,
                    size: file.size,
                    type: file.type
                }));
                window.parent.postMessage({
                    type: 'folder-upload',
                    files: fileData
                }, '*');
            });
        </script>
    """
    st.components.v1.html(folder_upload_html, height=50)

def initialize_session_state():
    if 'manager' not in st.session_state:
        st.session_state.manager = None
    if 'system_ready' not in st.session_state:
        st.session_state.system_ready = False
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'api_key' not in st.session_state:
        st.session_state.api_key = os.getenv('GOOGLE_API_KEY', '')
    if 'folder_files' not in st.session_state:
        st.session_state.folder_files = []
    if 'folder_uploaded' not in st.session_state:
        st.session_state.folder_uploaded = False

def main():
    st.set_page_config(page_title="FolderFlow QA Assistant", page_icon="📚")
    st.markdown("<h1>FolderFlow for <span style='color: #D00F22; font-weight: bold; font-size: 50px;'>HILTI</span> Technologies</h1>", unsafe_allow_html=True)

    initialize_session_state()
    
    st.sidebar.markdown("<h3 style='color: white;'>Configuration</h3>", unsafe_allow_html=True)

    api_key = st.sidebar.text_input(
        label="Enter Google API Key",
        value=st.session_state.api_key,
        type="password",
        key="api_key_input",
        label_visibility="visible",
        help="Enter your Google API key here"
    )

    st.markdown(
        """
        <style>
        [data-testid="stTextInput"] label {
            color: white !important;
        }
        [data-testid="stFileUploader"] label {
            color: white !important;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    if api_key:
        st.session_state.api_key = api_key
        os.environ["GOOGLE_API_KEY"] = api_key
        genai.configure(api_key=api_key)
    
    st.sidebar.markdown("<h3 style='color: white;'>Upload Documents</h3>", unsafe_allow_html=True)
    
    # Add folder upload button
    add_folder_upload_button()
    
    # Regular file upload
    uploaded_files = st.sidebar.file_uploader(
        "Or upload individual files",
        type=['pdf', 'docx', 'pptx', 'ppt'],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files",
        label_visibility="visible"
    )

    # Handle folder upload
    if st.session_state.get('folder_uploaded', False):
        with st.spinner("Processing folder..."):
            manager = DocumentManager()
            if manager.setup_qa_system(st.session_state.folder_files, is_folder=True):
                st.session_state.manager = manager
                st.session_state.system_ready = True
                st.sidebar.success("✅ Folder processed successfully!")
            else:
                st.sidebar.error("❌ Failed to process folder.")
            st.session_state.folder_files = []
            st.session_state.folder_uploaded = False

    # Handle individual file upload
    if uploaded_files and api_key and st.sidebar.button("Process Documents", type="primary"):
        with st.spinner("Processing documents..."):
            manager = DocumentManager()
            if manager.setup_qa_system(uploaded_files):
                st.session_state.manager = manager
                st.session_state.system_ready = True
                st.sidebar.markdown("""
                    <div style='background-color: #1aff1a; 
                                color: #000000;
                                padding: 10px;
                                border-radius: 5px;
                                font-weight: bold;'>
                        ✅ Documents processed successfully!
                    </div>
                    """, 
                    unsafe_allow_html=True
                )
            else:
                st.sidebar.markdown("""
                    <div style='background-color: #ffcccc; 
                                color: #cc0000;
                                padding: 10px;
                                border-radius: 5px;
                                font-weight: bold;'>
                        ❌ Setup failed. Please try again.
                    </div>
                    """, 
                    unsafe_allow_html=True
                )

    # Display processed files
    if st.session_state.system_ready and st.session_state.manager.processed_files:
        st.sidebar.header("Processed Files")
        for file in st.session_state.manager.processed_files:
            st.sidebar.text(f"✓ {file}")

    # Chat interface
    if st.session_state.system_ready:
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])
                if "sources" in message:
                    st.markdown("**Sources:**")
                    for source in message["sources"]:
                        st.markdown(f"- {source}")

        if prompt := st.chat_input("Ask a question about your documents"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = st.session_state.manager.ask_question(prompt)
                    if "error" in response:
                        st.error(response["error"])
                    else:
                        st.write(response["answer"])
                        if response["sources"]:
                            st.markdown("**Sources:**")
                            for source in response["sources"]:
                                st.markdown(f"- {source}")
                        
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": response["answer"],
                            "sources": response["sources"]
                        })
    else:
        st.info("👋 Hi, Welcome to our prototype FolderFlow - More than just 'search' \n\n"
                "Prepared by Dibyanshu and Sajjad - Happy Exploring!\n\n"
                "Kindly follow these steps to begin:\n"
                "1. Enter your Google API key in the sidebar\n"
                "2. Upload your documents or folder\n"
                "3. Click 'Process Documents' to start")

if __name__ == "__main__":
    main()
